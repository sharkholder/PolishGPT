import os
import docx
import PyPDF2
import openpyxl
import logging
from pptx import Presentation
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
from PIL import Image, ImageOps
import pytesseract
from openai import OpenAI  # Import OpenAI client

# Initialize the Flask application
app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'docx', 'pdf', 'xlsx', 'pptx', 'png', 'jpg', 'jpeg'}
file_cache = {}  # Cache for storing uploaded files' content

# Ensure the upload folder exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Configure the app
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max upload size

# Initialize OpenAI client with API key
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY", "Your_Key_Here"))

def allowed_file(filename):
    """Check if the file has an allowed extension."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def read_docx(file_path):
    """Read a DOCX file and return its text."""
    try:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)
    except Exception as e:
        logging.error(f"Failed to read DOCX: {e}")
        return str(e)

def read_pdf(file_path):
    """Read a PDF file and return its text."""
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            return '\n'.join(page.extract_text() for page in reader.pages if page.extract_text())
    except Exception as e:
        logging.error(f"Failed to read PDF: {e}")
        return str(e)

def read_excel(file_path):
    """Read an Excel file and return its content as text."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet = wb.active
        return '\n'.join(
            ','.join(str(cell) if cell is not None else '' for cell in row)
            for row in sheet.iter_rows(values_only=True)
        )
    except Exception as e:
        logging.error(f"Failed to read Excel: {e}")
        return str(e)

def extract_text_from_image(image_path):
    """Extract text from an image using OCR (Optical Character Recognition)."""
    try:
        image = Image.open(image_path)
        image = ImageOps.grayscale(image)  # Convert to grayscale to improve OCR accuracy
        return pytesseract.image_to_string(image)
    except Exception as e:
        logging.error(f"Failed to read image: {e}")
        return "Error reading image"

def convert_ppt_to_text(file_path):
    """Convert PPTX slides to text."""
    try:
        presentation = Presentation(file_path)
        slide_texts = []
        for slide in presentation.slides:
            slide_text = ' '.join(shape.text for shape in slide.shapes if shape.has_text_frame)
            slide_texts.append(slide_text)
        return '\n'.join(slide_texts)
    except Exception as e:
        logging.error(f"Failed to read PPTX: {e}")
        return str(e)

@app.route('/', methods=['GET'])
def index():
    """Render the main page with the file upload form and question input."""
    return render_template('index.html')

@app.route('/upload_multiple', methods=['POST'])
def upload_multiple_files():
    """Handle multiple file upload requests."""
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "No files part"}), 400

    user_id = request.remote_addr
    if user_id not in file_cache:
        file_cache[user_id] = []

    uploaded_files_info = []

    for i, file in enumerate(files):
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            file_content = None
            if filename.endswith('.docx'):
                file_content = read_docx(file_path)
                file_type = 'DOCX'
            elif filename.endswith('.pdf'):
                file_content = read_pdf(file_path)
                file_type = 'PDF'
            elif filename.endswith('.xlsx'):
                file_content = read_excel(file_path)
                file_type = 'Excel'
            elif filename.endswith('.pptx'):
                file_content = convert_ppt_to_text(file_path)
                file_type = 'PPTX'
            elif filename.lower().endswith(('png', 'jpg', 'jpeg')):
                file_content = extract_text_from_image(file_path)
                file_type = 'Image'
            else:
                continue  # Skip unsupported file types

            # Handle reading errors
            if isinstance(file_content, str) and file_content.startswith("Error"):
                uploaded_files_info.append(
                    {"filename": filename, "status": f"Failed: {file_content}"})
            else:
                # Add this file's content to the cache
                file_cache[user_id].append(
                    {"filename": filename, "content": file_content, "file_type": file_type, "id": f"{filename}-{i}"})
                uploaded_files_info.append(
                    {"filename": filename, "status": "Uploaded successfully", "id": f"{filename}-{i}"})

            # Remove the file after processing
            os.remove(file_path)
        else:
            uploaded_files_info.append({"filename": file.filename, "status": "Unsupported file type"})

    return jsonify(uploaded_files_info), 200

@app.route('/ask', methods=['POST'])
def ask_question():
    """Receives a user's question and returns an answer based on the uploaded document."""
    question = request.form.get('question')
    user_id = request.remote_addr
    if user_id in file_cache and file_cache[user_id]:
        document_texts = " ".join([file['content'] for file in file_cache[user_id]])
        try:
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "Below is an analysis of uploaded documents."},
                    {"role": "user", "content": document_texts},
                    {"role": "user", "content": question}
                ]
            )
            # Correct access to message content properties
            if response.choices and response.choices[0].message.content:
                return jsonify({'answer': response.choices[0].message.content})
            return jsonify({'error': "No response from AI."})
        except Exception as e:
            logging.error(f"Error with OpenAI service: {e}")
            return jsonify({'error': str(e)}), 500
    else:
        return jsonify({'error': "No documents uploaded or documents do not contain text."}), 400


@app.route('/view/text', methods=['GET'])
def view_all_texts():
    """Display all uploaded text files' content in order."""
    user_id = request.remote_addr
    files = file_cache.get(user_id, [])
    if not files:
        return jsonify(
            {"files": [{"filename": "No files", "content": "No content available.", "file_type": "Unknown"}]})

    # Group files by type
    categorized_files = {"DOCX": [], "PDF": [], "Excel": [], "PPTX": [], "Image": [], "Unknown": []}
    for file in files:
        categorized_files[file.get("file_type", "Unknown")].append(file)

    return jsonify({"categorized_files": categorized_files})

@app.route('/cancel/<filename>', methods=['DELETE'])
def cancel_file_upload(filename):
    """Remove a cancelled file from the cache."""
    user_id = request.remote_addr
    files = file_cache.get(user_id, [])

    # Match by filename and remove
    remaining_files = [file for file in files if file['id'] != filename]
    file_cache[user_id] = remaining_files

    # Return result
    if len(remaining_files) < len(files):
        return jsonify({"success": True})
    else:
        return jsonify({"success": False, "error": "File not found"}), 404

# Start the application
if __name__ == '__main__':
    app.run(debug=True)

